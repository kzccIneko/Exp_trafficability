#!/usr/bin/env python3
"""
组会汇报PPT生成脚本
基于方向敏感越野通行能力模型研究
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import numpy as np
from pathlib import Path
import io
import os

# 设置中文字体
plt.rcParams['font.sans-serif'] = ['SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

class PPTGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.slide_width = Inches(13.333)
        self.slide_height = Inches(7.5)
        
        # 颜色方案
        self.colors = {
            'primary': RGBColor(0, 102, 204),      # 主色：蓝色
            'secondary': RGBColor(255, 102, 0),     # 辅色：橙色
            'accent': RGBColor(0, 153, 51),         # 强调：绿色
            'dark': RGBColor(51, 51, 51),           # 深色
            'light': RGBColor(240, 240, 240),       # 浅色
            'white': RGBColor(255, 255, 255),
            'gray': RGBColor(128, 128, 128)
        }
        
        # 设置幻灯片尺寸
        self.prs.slide_width = self.slide_width
        self.prs.slide_height = self.slide_height
        
        # 文件路径
        self.base_dir = Path("D:/VSCode Program/通行能力分析_研")
        self.gpt_output_dir = self.base_dir / "gpt_output" / "prototype_v4_2_full" / "outputs_real_1000"
        self.our_output_dir = self.base_dir / "prototype"
        
    def add_title_slide(self, title, subtitle, author, date):
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 标题布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # 设置副标题
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        
        # 添加作者和日期
        left = Inches(1)
        top = Inches(5.5)
        width = Inches(11)
        height = Inches(1)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{author}\n{date}"
        p.font.size = Pt(14)
        p.font.color.rgb = self.colors['gray']
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(self, title, content_text, notes_text=""):
        """添加内容幻灯片"""
        slide_layout = self.prs.slide_layouts[1]  # 标题和内容布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # 设置内容
        content_shape = slide.placeholders[1]
        content_shape.text = content_text
        content_shape.text_frame.paragraphs[0].font.size = Pt(18)
        
        # 添加讲者备注
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        
        return slide
    
    def add_image_slide(self, title, image_path, caption="", notes_text=""):
        """添加图片幻灯片"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.333)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['primary']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 添加图片
        if Path(image_path).exists():
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(12.333)
            height = Inches(5.5)
            
            pic = slide.shapes.add_picture(str(image_path), left, top, width, height)
            
            # 添加说明文字
            if caption:
                left = Inches(0.5)
                top = Inches(6.8)
                width = Inches(12.333)
                height = Inches(0.5)
                
                caption_box = slide.shapes.add_textbox(left, top, width, height)
                caption_tf = caption_box.text_frame
                caption_p = caption_tf.paragraphs[0]
                caption_p.text = caption
                caption_p.font.size = Pt(12)
                caption_p.font.color.rgb = self.colors['gray']
                caption_p.alignment = PP_ALIGN.CENTER
        
        # 添加讲者备注
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        
        return slide
    
    def add_formula_slide(self, title, formula_image_path, explanation_text, notes_text=""):
        """添加公式幻灯片"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.333)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['primary']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 添加公式图片
        if Path(formula_image_path).exists():
            left = Inches(1)
            top = Inches(1.2)
            width = Inches(11.333)
            height = Inches(3)
            
            pic = slide.shapes.add_picture(str(formula_image_path), left, top, width, height)
        
        # 添加解释文字
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(12.333)
        height = Inches(2.5)
        
        explanation_box = slide.shapes.add_textbox(left, top, width, height)
        explanation_tf = explanation_box.text_frame
        explanation_p = explanation_tf.paragraphs[0]
        explanation_p.text = explanation_text
        explanation_p.font.size = Pt(16)
        explanation_p.font.color.rgb = self.colors['dark']
        
        # 添加讲者备注
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        
        return slide
    
    def add_comparison_slide(self, title, image_paths, captions, notes_text=""):
        """添加对比幻灯片"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.333)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['primary']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 添加多张图片（横向排列）
        num_images = len(image_paths)
        if num_images > 0:
            img_width = (12.333 - 0.5 * (num_images - 1)) / num_images
            left_start = 0.5
            
            for i, (img_path, caption) in enumerate(zip(image_paths, captions)):
                if Path(img_path).exists():
                    left = Inches(left_start + i * (img_width + 0.5))
                    top = Inches(1.2)
                    width = Inches(img_width)
                    height = Inches(5)
                    
                    pic = slide.shapes.add_picture(str(img_path), left, top, width, height)
                    
                    # 添加说明文字
                    caption_box = slide.shapes.add_textbox(left, top + height + Inches(0.1), width, Inches(0.5))
                    caption_tf = caption_box.text_frame
                    caption_p = caption_tf.paragraphs[0]
                    caption_p.text = caption
                    caption_p.font.size = Pt(12)
                    caption_p.font.color.rgb = self.colors['gray']
                    caption_p.alignment = PP_ALIGN.CENTER
        
        # 添加讲者备注
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        
        return slide
    
    def add_table_slide(self, title, headers, data, notes_text=""):
        """添加表格幻灯片"""
        slide_layout = self.prs.slide_layouts[5]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.333)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['primary']
        title_p.alignment = PP_ALIGN.CENTER
        
        # 添加表格
        rows = len(data) + 1  # 包括表头
        cols = len(headers)
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11.333)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 设置表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            cell.text_frame.paragraphs[0].font.color.rgb = self.colors['white']
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 设置数据行
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 交替行颜色
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors['light']
        
        # 添加讲者备注
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text
        
        return slide
    
    def generate_formula_images(self):
        """生成关键公式的PNG图片"""
        formulas_dir = self.base_dir / "ppt_formulas"
        formulas_dir.mkdir(exist_ok=True)
        
        formulas = {
            "gradient_decomposition": r"$\nabla z = (p, q)$\n$g_\parallel = p\cos\theta + q\sin\theta$\n$g_\perp = -p\sin\theta + q\cos\theta$",
            "asymmetric_impedance": r"$R_\parallel(g_\parallel) = g_\parallel^2 / \tan^2\alpha_u$ (上坡)\n$R_\parallel(g_\parallel) = 0$ (缓下坡)\n$R_\parallel(g_\parallel) = g_\parallel^2 / \tan^2\alpha_d$ (陡下坡)",
            "cross_slope_impedance": r"$R_\perp(g_\perp) = 1 - \exp\left(-\left(\frac{|g_\perp|}{\tan\alpha_r}\right)^2\right)$",
            "non_compensatory": r"$R_s = 1 - (1 - R_\parallel)(1 - R_\perp)$",
            "unit_cost": r"$c_{\text{unit}} = -\ln(P + \varepsilon)$\n$P = 1 - R_s$",
            "anisotropy_index": r"$AI_p = \frac{Q_{90}(c_{\text{unit}}) - Q_{10}(c_{\text{unit}})}{Q_{90}(c_{\text{unit}}) + Q_{10}(c_{\text{unit}}) + \varepsilon}$"
        }
        
        for name, formula_text in formulas.items():
            fig, ax = plt.subplots(figsize=(10, 3))
            ax.text(0.5, 0.5, formula_text, 
                    fontsize=24, 
                    ha='center', va='center',
                    transform=ax.transAxes,
                    bbox=dict(boxstyle="round,pad=0.5", 
                             facecolor="white", 
                             edgecolor="lightgray"))
            ax.axis('off')
            
            # 保存为PNG
            formula_path = formulas_dir / f"{name}.png"
            plt.savefig(formula_path, dpi=150, bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            plt.close()
        
        return formulas_dir
    
    def create_presentation(self):
        """创建完整的演示文稿"""
        print("开始生成组会汇报PPT...")
        
        # 1. 标题幻灯片
        print("1. 创建标题幻灯片...")
        self.add_title_slide(
            "方向敏感的越野通行能力建模方法研究",
            "基于地形梯度分解的代价函数构建与验证",
            "NEKO",
            "2026年6月 组会汇报"
        )
        
        # 2. 目录幻灯片
        print("2. 创建目录幻灯片...")
        self.add_content_slide(
            "汇报目录",
            "1. 研究背景与问题\n"
            "2. 理论方法与创新\n"
            "   2.1 地形梯度分解\n"
            "   2.2 三段式非对称阻抗函数\n"
            "   2.3 横坡阻抗函数\n"
            "   2.4 非补偿性综合\n"
            "   2.5 方向敏感的代价模型\n"
            "3. 实验设计与验证\n"
            "   3.1 对比实验设计\n"
            "   3.2 定量验证指标\n"
            "   3.3 实验结果分析\n"
            "4. 参数敏感性分析\n"
            "5. 结论与展望",
            "今天汇报主要围绕老师上次提出的四个核心问题：\n"
            "1. 坡度不只是阻碍还可能是助力\n"
            "2. 坡度代价与行进方向相关\n"
            "3. 方法太简单缺乏理论深度\n"
            "4. 缺乏定量验证\n"
            "我们将通过完整的理论推导和实验验证来回答这些问题。"
        )
        
        # 3. 研究背景与问题
        print("3. 创建研究背景与问题幻灯片...")
        self.add_content_slide(
            "研究背景与核心问题",
            "• 越野通行能力建模：地理环境要素 → 通行阻碍度 → 综合通行能力\n"
            "• 传统方法：标量场模型（B0）\n"
            "  - 仅考虑坡度大小，忽略方向差异\n"
            "  - 无法区分上坡、下坡、横坡的不同影响\n"
            "• 老师提出的四个核心问题：\n"
            "  1. 坡度不只是阻碍还可能是助力\n"
            "  2. 坡度代价与行进方向密切相关\n"
            "  3. 现有方法理论深度不足\n"
            "  4. 缺乏客观定量验证方法",
            "老师的核心关切：我们需要一个有理论深度、能物理验证、方向敏感的通行能力模型。\n"
            "现有标量场模型无法捕捉地形的方向性特征，这正是我们研究的出发点。"
        )
        
        # 4. 理论方法 - 梯度分解
        print("4. 创建理论方法幻灯片...")
        formulas_dir = self.generate_formula_images()
        
        self.add_formula_slide(
            "2.1 地形梯度分解",
            formulas_dir / "gradient_decomposition.png",
            "• ∇z = (p, q)：地形表面梯度向量\n"
            "• p = ∂z/∂x：x方向坡度分量\n"
            "• q = ∂z/∂y：y方向坡度分量\n"
            "• g∥：沿车辆方向θ的纵坡分量\n"
            "• g⊥：垂直车辆方向θ的横坡分量\n"
            "• 物理意义：将地形坡度分解为车辆坐标系下的两个正交分量",
            "梯度分解是方向敏感建模的基础。通过将地形梯度∇z投影到车辆坐标系，我们得到：\n"
            "1. 纵坡g∥：直接影响车辆纵向动力学（上坡阻力/下坡助力）\n"
            "2. 横坡g⊥：影响车辆横向稳定性（侧翻风险）\n"
            "这回答了老师的第一个问题：坡度确实可能提供助力（下坡时），而不仅仅是阻碍。"
        )
        
        # 5. 三段式非对称阻抗
        self.add_formula_slide(
            "2.2 三段式非对称阻抗函数",
            formulas_dir / "asymmetric_impedance.png",
            "• R∥(g∥)：纵坡阻抗函数，分三段：\n"
            "  - 上坡(g∥>0)：R∥ = g∥²/tan²(αu)\n"
            "  - 缓下坡(|g∥|≤tan(αm))：R∥ = 0\n"
            "  - 陡下坡(g∥<-tan(αm))：R∥ = g∥²/tan²(αd)\n"
            "• αu：上坡临界角（15°）\n"
            "• αm：缓下坡死区角（5°）\n"
            "• αd：陡下坡临界角（15°）\n"
            "• 设计原理：反映车辆上坡阻力大、下坡有一定助力但过陡则失控的物理特性",
            "三段式设计体现了坡度的双重作用：\n"
            "1. 上坡段：坡度越大阻力越大，阻抗随坡度平方增长\n"
            "2. 缓下坡段：坡度提供助力，阻抗为0（死区设计）\n"
            "3. 陡下坡段：坡度过大导致失控风险，阻抗再次增大\n"
            "这直接回答了老师的问题：坡度确实可能是助力（缓下坡段），但需要分段建模。"
        )
        
        # 6. 横坡阻抗
        self.add_formula_slide(
            "2.3 横坡阻抗函数",
            formulas_dir / "cross_slope_impedance.png",
            "• R⊥(g⊥)：横坡阻抗函数，对称设计\n"
            "• 公式：R⊥ = 1 - exp(-(|g⊥|/tan(αr))²)\n"
            "• αr：横坡临界角（15°）\n"
            "• 设计原理：\n"
            "  - 对称性：左右横坡风险相同\n"
            "  - 高斯形状：小横坡风险低，大横坡风险急剧增加\n"
            "  - 物理意义：横坡主要影响横向稳定性，与纵坡的物理机制不同",
            "横坡阻抗函数采用对称设计，因为：\n"
            "1. 左倾和右倾的侧翻风险理论上相同\n"
            "2. 横坡主要影响横向稳定性，而非纵向动力学\n"
            "3. 采用指数形式确保平滑过渡，避免阻抗突变"
        )
        
        # 7. 非补偿性综合
        self.add_formula_slide(
            "2.4 非补偿性综合",
            formulas_dir / "non_compensatory.png",
            "• Rs：综合通行阻抗\n"
            "• 公式：Rs = 1 - (1 - R∥)(1 - R⊥)\n"
            "• 非补偿性：任一维度的高风险不会被另一维度的低风险补偿\n"
            "• 设计原理：\n"
            "  - 如果R∥=0.9（上坡很陡），即使R⊥=0.0（完全平坦），Rs=0.9\n"
            "  - 如果R⊥=0.9（横坡很陡），即使R∥=0.0（完全平坦），Rs=0.9\n"
            "  - 体现'短板效应'：通行能力由最危险维度决定",
            "非补偿性综合是本方法的重要创新点：\n"
            "1. 传统加权求和是补偿性的：一个维度的低风险可以补偿另一个维度的高风险\n"
            "2. 实际通行中：上坡很陡时，即使路面平坦也无法通行\n"
            "3. 横坡很陡时，即使纵坡平坦也容易侧翻\n"
            "4. 非补偿性更符合实际物理限制"
        )
        
        # 8. 单位距离代价
        self.add_formula_slide(
            "2.5 方向敏感的代价模型",
            formulas_dir / "unit_cost.png",
            "• cunit：单位距离代价（每米通行成本）\n"
            "• P = 1 - Rs：通行概率\n"
            "• 公式：cunit = -ln(P + ε)\n"
            "• ε：小常数（1e-10），避免ln(0)\n"
            "• 设计原理：\n"
            "  - 概率→代价的转换：P越大（越安全），cunit越小\n"
            "  - 对数转换：确保代价为正，且对低概率更敏感\n"
            "  - 单位距离：便于A*算法计算路径总代价",
            "单位距离代价模型将通行概率转换为可加性代价：\n"
            "1. 通行概率P ∈ [0,1]：P=1表示完全安全，P=0表示完全不可通行\n"
            "2. 对数转换：cunit = -ln(P)，当P=1时cunit=0，P→0时cunit→∞\n"
            "3. 这种转换使得路径总代价可以简单累加：J = Σ cunit × Δs\n"
            "4. 方向敏感：对于同一栅格，不同方向θ的cunit不同"
        )
        
        # 9. 实验设计
        print("5. 创建实验设计幻灯片...")
        self.add_content_slide(
            "3. 实验设计与验证",
            "• 三模型对比：\n"
            "  - B0：标量坡度模型（仅坡度大小）\n"
            "  - B1：仅纵坡模型（方向敏感，但忽略横坡）\n"
            "  - Ours：纵坡+横坡模型（完整方向敏感模型）\n"
            "• 实验数据：雅砻江真实DEM（1000×1000，分辨率28.8m）\n"
            "• 验证方法：\n"
            "  - 梯度流线验证（物理验证）\n"
            "  - 定量指标验证\n"
            "  - 参数敏感性分析",
            "实验设计的核心思想：\n"
            "1. 三模型对比：验证横坡信息的重要性\n"
            "2. 真实DEM：确保实验结果的实际意义\n"
            "3. 多维度验证：从物理合理性和定量指标两个角度验证\n"
            "这直接回应老师第四个问题：我们设计了客观定量的验证方法。"
        )
        
        # 10. 定量验证指标
        self.add_formula_slide(
            "3.2 定量验证指标",
            formulas_dir / "anisotropy_index.png",
            "• AUC：各向异性曲线下面积\n"
            "  - 衡量代价场的方向敏感程度\n"
            "  - AUC越大，方向敏感性越强\n"
            "• DC：方向一致性\n"
            "  - 最优方向与梯度流线方向的一致性\n"
            "  - DC越高，模型越符合物理规律\n"
            "• AIp：分位数各向异性指数\n"
            "  - 使用Q90和Q10，减弱极端值影响\n"
            "  - AIp越小，代价场越各向同性",
            "定量验证指标的设计原则：\n"
            "1. AUC：基于代价玫瑰图，量化不同方向代价差异\n"
            "2. DC：基于梯度流线（物理验证基准），验证模型是否符合地形约束\n"
            "3. AIp：使用分位数而非极值，提高指标鲁棒性\n"
            "这些指标为老师第四个问题提供了客观验证手段。"
        )
        
        # 11. 实验结果 - 场统计
        print("6. 创建实验结果幻灯片...")
        self.add_image_slide(
            "3.3.1 场统计指标对比",
            self.gpt_output_dir / "03_各向异性指数对比.png",
            "三模型场统计指标对比",
            "实验结果分析：\n"
            "1. B0（标量模型）：各向异性指数AIp=0，完全无方向敏感性\n"
            "2. B1（仅纵坡）：AIp=0.92，方向敏感性过强（可能夸大差异）\n"
            "3. Ours（纵坡+横坡）：AIp=0.17，适度的方向敏感性\n"
            "结论：完整模型提供了更平衡的方向敏感性"
        )
        
        # 12. 实验结果 - 代价场对比
        self.add_image_slide(
            "3.3.2 三模型最小代价场对比",
            self.gpt_output_dir / "02_三模型最小代价场对比.png",
            "B0、B1、Ours三模型最小代价场可视化对比",
            "代价场可视化对比：\n"
            "1. B0：代价分布均匀，无法区分不同地形特征\n"
            "2. B1：代价分布极端化，某些区域代价极低（接近0）\n"
            "3. Ours：代价分布更合理，反映真实通行难度\n"
            "这表明完整模型能更准确地刻画地形通行特性"
        )
        
        # 13. 实验结果 - 路径规划
        self.add_image_slide(
            "3.3.3 三模型路径规划对比",
            self.gpt_output_dir / "05_三模型路径规划对比.png",
            "三模型A*路径规划结果对比",
            "路径规划结果对比：\n"
            "1. B0路径：沿边界绕行，不合理\n"
            "2. B1路径：过于激进，穿越陡坡区域\n"
            "3. Ours路径：沿地形特征线走，符合物理规律\n"
            "关键指标：Ours的路径长度J=20802，与B0接近但路径更合理"
        )
        
        # 14. 实验结果 - 各向异性
        self.add_image_slide(
            "3.3.4 各向异性指数对比",
            self.gpt_output_dir / "03_各向异性指数对比.png",
            "B1与Ours各向异性指数空间分布对比",
            "各向异性指数分析：\n"
            "1. B1：AIp=0.92，空间分布不均匀\n"
            "2. Ours：AIp=0.17，空间分布更均匀\n"
            "3. 含义：Ours模型在不同区域都保持适度的方向敏感性\n"
            "这验证了非补偿性综合的有效性"
        )
        
        # 15. 参数敏感性分析
        print("7. 创建参数敏感性分析幻灯片...")
        self.add_content_slide(
            "4. 参数敏感性分析",
            "• 关键参数：\n"
            "  - αu：上坡临界角（5°-30°）\n"
            "  - αm：缓下坡死区角（0°-12°）\n"
            "  - αd：陡下坡临界角（5°-30°）\n"
            "  - αr：横坡临界角（5°-20°）\n"
            "• 敏感性分析方法：\n"
            "  - 固定其他参数，变化单一参数\n"
            "  - 观察对代价场和路径规划的影响\n"
            "• 目标：确定参数的合理范围，验证模型鲁棒性",
            "参数敏感性分析的目的：\n"
            "1. 确定参数的物理意义和合理范围\n"
            "2. 验证模型对参数变化的鲁棒性\n"
            "3. 为实际应用提供参数选择指导\n"
            "这是老师第三个问题（方法太简单）的重要回应：我们进行了系统的参数分析。"
        )
        
        # 16. 路径规划指标表格
        print("8. 创建路径规划指标表格幻灯片...")
        headers = ["模型", "路径总代价 J", "路径长度 L(m)", "平均代价 J/L", "绕行率 L/D"]
        data = [
            ["B0 标量模型", "26213", "56674", "0.463", "1.74"],
            ["B1 仅纵坡", "5433", "64383", "0.084", "1.98"],
            ["Ours 纵坡+横坡", "20802", "56771", "0.366", "1.74"]
        ]
        
        self.add_table_slide(
            "4.1 路径规划指标对比",
            headers,
            data,
            "路径规划指标分析：\n"
            "1. B1模型代价最低（5433），但绕行率最高（1.98），说明路径过于激进\n"
            "2. B0和Ours绕行率相近（1.74），但Ours代价更低（20802 vs 26213）\n"
            "3. Ours模型在保持合理路径的同时，降低了通行代价\n"
            "这验证了横坡信息对路径规划的重要性"
        )
        
        # 16. 结论与展望
        print("8. 创建结论与展望幻灯片...")
        self.add_content_slide(
            "5. 结论与展望",
            "• 主要贡献：\n"
            "  1. 提出方向敏感的通行能力建模方法\n"
            "  2. 设计三段式非对称阻抗函数，体现坡度双重作用\n"
            "  3. 引入非补偿性综合，避免风险补偿效应\n"
            "  4. 建立基于梯度流线的物理验证方法\n"
            "  5. 设计定量验证指标体系\n"
            "• 回应老师四个问题：\n"
            "  ✓ 坡度可能助力（缓下坡段设计）\n"
            "  ✓ 方向敏感建模（梯度分解）\n"
            "  ✓ 理论深度（三段式函数+非补偿性综合）\n"
            "  ✓ 定量验证（梯度流线+多指标验证）\n"
            "• 未来工作：\n"
            "  - 多尺度DEM融合\n"
            "  - 动态通行能力模型\n"
            "  - 实时路径规划算法",
            "总结与展望：\n"
            "1. 我们成功构建了方向敏感的通行能力模型\n"
            "2. 通过实验验证了方法的有效性\n"
            "3. 直接回应了老师的四个核心问题\n"
            "4. 为后续研究奠定了基础\n"
            "感谢老师的指导和建议！"
        )
        
        # 保存PPT
        output_path = self.base_dir / "组会汇报_20260601.pptx"
        self.prs.save(str(output_path))
        print(f"PPT生成完成：{output_path}")
        
        return output_path

if __name__ == "__main__":
    generator = PPTGenerator()
    output_path = generator.create_presentation()
    print(f"成功生成PPT文件：{output_path}")