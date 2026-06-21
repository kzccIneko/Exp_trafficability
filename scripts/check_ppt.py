#!/usr/bin/env python3
"""检查生成的PPT内容"""

from pptx import Presentation
from pathlib import Path

def check_ppt():
    ppt_path = Path("D:/VSCode Program/通行能力分析_研/组会汇报_20260601.pptx")
    
    if not ppt_path.exists():
        print("PPT文件不存在")
        return
    
    prs = Presentation(str(ppt_path))
    
    print(f"PPT文件大小：{ppt_path.stat().st_size / 1024 / 1024:.2f} MB")
    print(f"幻灯片数量：{len(prs.slides)}")
    print("\n幻灯片内容概览：")
    
    for i, slide in enumerate(prs.slides, 1):
        # 获取幻灯片标题
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        
        # 获取讲者备注
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text[:100] + "..." if slide.notes_slide.notes_text_frame.text else ""
        
        print(f"{i:2d}. {title[:50]:<50} | 备注：{notes[:30]}")
    
    # 检查是否有图片
    image_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image_count += 1
    
    print(f"\n图片总数：{image_count}")
    
    # 检查是否有表格
    table_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
    
    print(f"表格总数：{table_count}")

if __name__ == "__main__":
    check_ppt()